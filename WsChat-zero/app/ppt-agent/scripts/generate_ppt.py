#!/usr/bin/env python3
"""Generate a professional PowerPoint file from a JSON slide payload."""

import argparse
import json
import os
import sys
import random

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

STYLES = {
    "简约": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title": RGBColor(0x22, 0x22, 0x22),
        "body": RGBColor(0x55, 0x55, 0x55),
        "accent": RGBColor(0x2F, 0x80, 0xED),
        "accent2": RGBColor(0x5B, 0xA0, 0xF0),
        "font": "Microsoft YaHei",
    },
    "商务": {
        "bg": RGBColor(0xF4, 0xF7, 0xFA),
        "title": RGBColor(0x10, 0x24, 0x3E),
        "body": RGBColor(0x3E, 0x4C, 0x59),
        "accent": RGBColor(0x0B, 0x69, 0xA3),
        "accent2": RGBColor(0x2C, 0x8C, 0xC8),
        "font": "Microsoft YaHei",
    },
    "科技": {
        "bg": RGBColor(0x0B, 0x10, 0x1A),
        "title": RGBColor(0x64, 0xD2, 0xFF),
        "body": RGBColor(0xD9, 0xE2, 0xEC),
        "accent": RGBColor(0x2D, 0x9C, 0xDB),
        "accent2": RGBColor(0x56, 0xC8, 0xFF),
        "font": "Microsoft YaHei",
    },
    "创意": {
        "bg": RGBColor(0xFF, 0xF7, 0xED),
        "title": RGBColor(0x9A, 0x34, 0x12),
        "body": RGBColor(0x7C, 0x2D, 0x12),
        "accent": RGBColor(0xEA, 0x58, 0x0C),
        "accent2": RGBColor(0xF0, 0x8C, 0x4A),
        "font": "Microsoft YaHei",
    },
    "学术": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title": RGBColor(0x1D, 0x4E, 0x89),
        "body": RGBColor(0x33, 0x33, 0x33),
        "accent": RGBColor(0x90, 0x2F, 0x2F),
        "accent2": RGBColor(0xB0, 0x5A, 0x5A),
        "font": "Microsoft YaHei",
    },
}

def _ensure_slides_array(data):
    if isinstance(data, dict):
        for key in ("slides", "data", "result", "items"):
            if key in data and isinstance(data[key], list):
                return data[key]
        for v in data.values():
            if isinstance(v, list) and len(v) > 0 and isinstance(v[0], dict):
                return v
    if isinstance(data, list):
        return data
    return []

def add_bg(slide, style):
    """Fill entire slide background."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = style["bg"]
    bg.line.fill.background()
    sp = bg._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, bg._element)

def accent_bar(slide, style, left=0, top=0, width=0.08, height=7.5):
    """Vertical accent bar on the left."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = style["accent"]
    bar.line.fill.background()

def bottom_line(slide, style):
    """Thin line at the bottom."""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = style["accent"]
    line.line.fill.background()

def page_number(slide, num, total, style):
    """Page number at bottom right."""
    box = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    p.alignment = PP_ALIGN.RIGHT

def add_text_box(slide, left, top, width, height, text, size=18,
                 color=None, bold=False, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    f = p.font
    f.size = Pt(size)
    f.bold = bold
    f.name = font_name
    if color:
        f.color.rgb = color
    return box

def cover_slide(slide, title, subtitle, style, total):
    add_bg(slide, style)
    # Large accent block
    block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(13.333), Inches(3.2))
    block.fill.solid()
    block.fill.fore_color.rgb = style["accent"]
    block.line.fill.background()

    add_text_box(slide, 0.8, 1.0, 11.5, 1.8, title,
                 size=40, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, align=PP_ALIGN.LEFT)
    add_text_box(slide, 0.8, 3.8, 11.5, 2.0, subtitle[:200] if subtitle else "",
                 size=18, color=style["body"], align=PP_ALIGN.LEFT)
    page_number(slide, 1, total, style)

def section_slide(slide, title, style, total, num):
    add_bg(slide, style)
    accent_bar(slide, style, left=0.3, top=0, width=0.1, height=7.5)
    add_text_box(slide, 1.0, 2.5, 11.0, 2.0, title,
                 size=32, color=style["accent"], bold=True, align=PP_ALIGN.LEFT)
    page_number(slide, num, total, style)

def content_slide(slide, title, items, style, total, num):
    add_bg(slide, style)
    accent_bar(slide, style)

    # Title with underline accent
    add_text_box(slide, 0.8, 0.4, 11.5, 0.7, title,
                 size=26, color=style["title"], bold=True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(1.1), Inches(2.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = style["accent"]
    line.line.fill.background()

    # Content bullets in two columns if many items
    top = 1.5
    col_width = 5.8
    mid = len(items) // 2 + len(items) % 2
    for i, item in enumerate(items[:12]):
        col = 0.8 if i < mid else 6.8
        row = i if i < mid else i - mid
        y = top + row * 0.6
        bullet = f"  {item}" if item.startswith(("数", "数", "数", "数")) else f"▸ {item}"  # Keep it simple
        bullet = f"• {item}"
        add_text_box(slide, col, y, col_width, 0.55, bullet,
                     size=14, color=style["body"])
    bottom_line(slide, style)
    page_number(slide, num, total, style)

def image_slide(slide, title, desc, style, total, num):
    add_bg(slide, style)
    accent_bar(slide, style)
    add_text_box(slide, 0.8, 0.4, 11.5, 0.7, title,
                 size=26, color=style["title"], bold=True)
    # Image placeholder with styled border
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(1.5), Inches(10.3), Inches(4.8))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
    box.line.color.rgb = style["accent"]
    box.line.width = Pt(1.5)
    add_text_box(slide, 2.0, 3.0, 9.3, 1.5, f"[配图] {desc}",
                 size=16, color=style["accent2"], align=PP_ALIGN.CENTER)
    bottom_line(slide, style)
    page_number(slide, num, total, style)

def chart_slide(slide, title, chart_type, style, total, num):
    add_bg(slide, style)
    accent_bar(slide, style)
    add_text_box(slide, 0.8, 0.4, 11.5, 0.7, title,
                 size=26, color=style["title"], bold=True)
    # Chart placeholder
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.0))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    box.line.color.rgb = style["accent2"]
    box.line.width = Pt(1)
    add_text_box(slide, 2.0, 3.0, 9.3, 1.0, f"[{chart_type}图表]",
                 size=18, color=style["accent2"], align=PP_ALIGN.CENTER)
    bottom_line(slide, style)
    page_number(slide, num, total, style)

def generate_ppt(slides_data, style_name, output_path):
    slides_data = _ensure_slides_array(slides_data)
    if not slides_data:
        raise ValueError("No valid slides data found")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]
    style = STYLES.get(style_name, STYLES["简约"])
    total = len(slides_data)

    for idx, sd in enumerate(slides_data):
        slide = prs.slides.add_slide(layout)
        num = idx + 1
        title = sd.get("title", f"第 {num} 页")
        content = sd.get("content", "")
        need_img = bool(sd.get("needImage", False))
        img_desc = sd.get("imageDescription", "")
        chart = sd.get("chart", "")

        if num == 1:
            cover_slide(slide, title, content, style, total)
        elif chart:
            chart_slide(slide, title, chart, style, total, num)
        elif need_img and img_desc:
            image_slide(slide, title, img_desc, style, total, num)
        else:
            items = [line.strip().lstrip("•-▸").strip()
                     for line in content.split("\n") if line.strip()]
            if not items:
                items = [content.strip() or "待补充内容"]
            content_slide(slide, title, items, style, total, num)

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    return output_path

def main():
    parser = argparse.ArgumentParser(description="生成专业 PPTX 文件")
    parser.add_argument("--input", required=True, help="输入 JSON 文件路径")
    parser.add_argument("--output", required=True, help="输出 PPTX 文件路径")
    parser.add_argument("--style", default="简约", help="PPT 风格")
    args = parser.parse_args()

    with open(args.input, "r", encoding="utf-8-sig") as f:
        slides_data = json.load(f)

    result_path = generate_ppt(slides_data, args.style, args.output)
    sys.stdout.reconfigure(encoding="utf-8")
    print(json.dumps({"filePath": result_path}, ensure_ascii=False))

if __name__ == "__main__":
    main()
